"""Format adapters that create traceable Markdown without modifying source files."""

from __future__ import annotations

from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
import hashlib
import html
import json
import math
import multiprocessing
import os
from pathlib import Path
import re
import signal
import subprocess
import tempfile
import threading
import time
import unicodedata
from typing import Any, Callable, Iterable

import yaml

from company_wiki.canonical_ingest import IngestService, ParserResult
from company_wiki.parser_adapters import (
    PAGE_AWARE_PDF_PARSER_NAME,
    PageAwarePDFResult,
    adapt_pdf_pages,
)
from company_wiki.source_contract import (
    EvidenceCoordinates,
    ParseStatus,
    QualityFlag,
    SourceManifest,
)

from .admission import processing_priority_sql
from .artifact_handle import ARTIFACT_HANDLE_SCHEMA_VERSION
from .models import CatalogConfig, NORMALIZER_VERSION, ProcessingReport
from .store import CatalogStore, canonical_json


_NORMALIZER_NAME = "source_catalog_normalizer"
_SENTENCE_BREAK_RE = re.compile(r"\n\s*\n+")
DOCUMENT_PARSE_TIMEOUT_CODE = "document_parse_timeout"
_UNSUPPORTED_CHILD_ERROR_TYPES = {
    "BadZipFile",
    "EmptyFileError",
    "FileDataError",
    "InvalidFileException",
    "PackageNotFoundError",
    "UnsupportedDocumentError",
    "XLRDError",
}
_UNSUPPORTED_VALUE_ERROR_MESSAGES = {
    "PDF is password protected",
    "XLS contains no non-empty sheets",
    "XLSX contains no non-empty sheets",
}


def _utc_iso(epoch: float | None = None) -> str:
    """UTC timestamp as a second-precision ISO-8601 string (sortable as text)."""
    import time as _time
    from datetime import datetime, timezone

    if epoch is None:
        epoch = _time.time()
    return datetime.fromtimestamp(epoch, tz=timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


@dataclass(frozen=True)
class _Normalized:
    body: str
    parser_results: tuple[ParserResult, ...]
    parser_name: str
    parser_version: str
    status: str
    quality_flags: tuple[str, ...]
    error: str | None = None
    # ZR-501: broker_research metadata contract — page count known by the
    # page-aware parser (None when unknown; never invented).
    page_count: int | None = None
    # ZR-502: raw first-page text for homepage identity verification.
    first_page_text: str | None = None


class NormalizationProcessError(RuntimeError):
    """Base class for failures in the isolated document parser process."""


class NormalizationCancelledError(NormalizationProcessError):
    """Raised when a stop request cancels an active parser process."""


class NormalizationTimeoutError(NormalizationProcessError):
    """Raised when one document exceeds its independent parser deadline."""


class ParserProcessError(NormalizationProcessError):
    """Raised when the parser child fails or returns an invalid envelope."""


class UnsupportedDocumentError(ParserProcessError):
    """Raised when a parser proves that the source bytes are malformed."""


class ParserResultTooLargeError(ParserProcessError):
    """Raised when parser output exceeds the configured IPC size limit."""


class ParserResultProtocolError(ParserProcessError):
    """Raised when parser output violates the normalized-result contract."""


def _normalized_to_payload(
    normalized: Any, *, expected_source_id: str
) -> dict[str, Any]:
    if not isinstance(normalized, _Normalized):
        raise ParserResultProtocolError("parser must return _Normalized")
    parser_results: list[dict[str, Any]] = []
    for result in normalized.parser_results:
        if not isinstance(result, ParserResult):
            raise ParserResultProtocolError(
                "normalized parser_results must contain ParserResult values"
            )
        if result.source_id != expected_source_id:
            raise ParserResultProtocolError(
                "parser result source_id does not match the input manifest"
            )
        span_payload = result.to_evidence_span().to_dict()
        parser_results.append(
            {
                "source_id": span_payload["source_id"],
                "coordinates": span_payload["coordinates"],
                "raw_text": span_payload["raw_text"],
                "structured_value": span_payload["structured_value"],
                "parser_name": span_payload["parser_name"],
                "parser_version": span_payload["parser_version"],
                "parse_status": span_payload["parse_status"],
                "quality_flags": span_payload["quality_flags"],
            }
        )
    return {
        "schema_version": "1.0",
        "source_id": expected_source_id,
        "body": normalized.body,
        "parser_results": parser_results,
        "parser_name": normalized.parser_name,
        "parser_version": normalized.parser_version,
        "status": normalized.status,
        "quality_flags": list(normalized.quality_flags),
        "error": normalized.error,
        # ZR-501: page count travels the isolated-parser envelope so the
        # normalized artifact frontmatter carries it end-to-end.
        "page_count": normalized.page_count,
        # ZR-502: first-page text travels too (homepage identity check).
        "first_page_text": normalized.first_page_text,
    }


def _normalized_from_payload(payload: Any, *, expected_source_id: str) -> _Normalized:
    if not isinstance(payload, dict):
        raise ParserResultProtocolError("parser result payload must be an object")
    expected_fields = {
        "schema_version",
        "source_id",
        "body",
        "parser_results",
        "parser_name",
        "parser_version",
        "status",
        "quality_flags",
        "error",
        "page_count",
        "first_page_text",
    }
    if set(payload) != expected_fields or payload.get("schema_version") != "1.0":
        raise ParserResultProtocolError("parser result payload schema is invalid")
    if payload.get("source_id") != expected_source_id:
        raise ParserResultProtocolError(
            "parser result payload source_id does not match the input manifest"
        )
    for field in ("body", "parser_name", "parser_version", "status"):
        if not isinstance(payload[field], str):
            raise ParserResultProtocolError(f"{field} must be text")
    if payload["status"] not in {"completed", "partial", "unsupported", "failed"}:
        raise ParserResultProtocolError("status is invalid")
    if not isinstance(payload["quality_flags"], list) or any(
        not isinstance(item, str) for item in payload["quality_flags"]
    ):
        raise ParserResultProtocolError("quality_flags must be an array of text")
    if payload["error"] is not None and not isinstance(payload["error"], str):
        raise ParserResultProtocolError("error must be text or null")
    page_count = payload.get("page_count")
    if page_count is not None and (
        isinstance(page_count, bool)
        or not isinstance(page_count, int)
        or page_count < 1
    ):
        raise ParserResultProtocolError("page_count must be a positive integer or null")
    first_page_text = payload.get("first_page_text")
    if first_page_text is not None and not isinstance(first_page_text, str):
        raise ParserResultProtocolError("first_page_text must be text or null")
    raw_results = payload.get("parser_results")
    if not isinstance(raw_results, list):
        raise ParserResultProtocolError("parser_results must be an array")
    parser_results: list[ParserResult] = []
    for item in raw_results:
        if not isinstance(item, dict):
            raise ParserResultProtocolError("parser result item must be an object")
        try:
            parser_result = ParserResult(
                source_id=item["source_id"],
                coordinates=EvidenceCoordinates.from_dict(item["coordinates"]),
                raw_text=item["raw_text"],
                structured_value=item["structured_value"],
                parser_name=item["parser_name"],
                parser_version=item["parser_version"],
                parse_status=item["parse_status"],
                quality_flags=item["quality_flags"],
            )
        except (KeyError, TypeError, ValueError) as exc:
            raise ParserResultProtocolError(
                f"invalid parser result item: {type(exc).__name__}: {str(exc)[:500]}"
            ) from exc
        if parser_result.source_id != expected_source_id:
            raise ParserResultProtocolError(
                "parser result source_id does not match the input manifest"
            )
        parser_results.append(parser_result)
    try:
        return _Normalized(
            body=payload["body"],
            parser_results=tuple(parser_results),
            parser_name=payload["parser_name"],
            parser_version=payload["parser_version"],
            status=payload["status"],
            quality_flags=tuple(payload["quality_flags"]),
            error=payload["error"],
            page_count=page_count,
            first_page_text=first_page_text,
        )
    except (TypeError, ValueError) as exc:
        raise ParserResultProtocolError(
            f"invalid normalized result: {type(exc).__name__}: {str(exc)[:500]}"
        ) from exc


def _write_parser_envelope(
    result_path: Path, envelope: dict[str, Any], *, result_max_bytes: int
) -> None:
    encoded = json.dumps(
        envelope,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    ).encode("utf-8")
    if envelope.get("ok") and len(encoded) > result_max_bytes:
        encoded = json.dumps(
            {
                "schema_version": "1.0",
                "ok": False,
                "error_type": "ParserResultTooLargeError",
                "error": (
                    f"parser result is {len(encoded)} bytes; limit is "
                    f"{result_max_bytes} bytes"
                ),
            },
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    temporary = result_path.with_name(result_path.name + f".{os.getpid()}.tmp")
    try:
        temporary.write_bytes(encoded)
        os.replace(temporary, result_path)
    finally:
        temporary.unlink(missing_ok=True)


def _terminate_windows_process_tree(pid: int) -> None:
    try:
        subprocess.run(
            ["taskkill.exe", "/PID", str(pid), "/T", "/F"],
            stdin=subprocess.DEVNULL,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            creationflags=subprocess.CREATE_NO_WINDOW,
            timeout=10,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired):
        pass


def _exit_if_parser_parent_dies(parent_liveness: Any) -> None:
    try:
        parent_liveness.recv()
    except (EOFError, OSError):
        if os.name == "nt":
            _terminate_windows_process_tree(os.getpid())
        else:
            # The parser called os.setsid() on startup, so it is the leader of
            # its own process group.  Kill the whole group so descendants
            # spawned by the parser are not left orphaned when the parent
            # monitor crashes.
            try:
                os.killpg(os.getpgid(0), signal.SIGKILL)
            except (OSError, ProcessLookupError):
                pass
        os._exit(70)


def _parser_process_entry(
    start_event: Any,
    result_path_text: str,
    path_text: str,
    manifest_payload: dict[str, Any],
    docling_path_text: str | None,
    result_max_bytes: int,
    parent_liveness: Any,
    parser: Callable[[Path, SourceManifest, Path | None], _Normalized],
) -> None:
    if os.name != "nt":
        try:
            os.setsid()
        except OSError:
            pass
    start_event.wait()
    monitor = threading.Thread(
        target=_exit_if_parser_parent_dies,
        args=(parent_liveness,),
        name="source-catalog-parser-parent-watch",
        daemon=True,
    )
    monitor.start()
    result_path = Path(result_path_text)
    try:
        manifest = SourceManifest.from_dict(manifest_payload)
        normalized = parser(
            Path(path_text),
            manifest,
            Path(docling_path_text) if docling_path_text is not None else None,
        )
        envelope = {
            "schema_version": "1.0",
            "ok": True,
            "result": _normalized_to_payload(
                normalized, expected_source_id=manifest.source_id
            ),
        }
    except BaseException as exc:  # noqa: BLE001 - error is returned to the parent
        envelope = {
            "schema_version": "1.0",
            "ok": False,
            "error_type": type(exc).__name__,
            "error": str(exc)[:1000],
        }
    _write_parser_envelope(result_path, envelope, result_max_bytes=result_max_bytes)


class _WindowsParserJob:
    """Kill a parser and its descendants if the owning worker handle closes."""

    def __init__(self, process: multiprocessing.Process):
        self._handle: int | None = None
        if os.name != "nt":
            return
        import ctypes
        from ctypes import wintypes

        class _IoCounters(ctypes.Structure):
            _fields_ = [
                ("ReadOperationCount", ctypes.c_ulonglong),
                ("WriteOperationCount", ctypes.c_ulonglong),
                ("OtherOperationCount", ctypes.c_ulonglong),
                ("ReadTransferCount", ctypes.c_ulonglong),
                ("WriteTransferCount", ctypes.c_ulonglong),
                ("OtherTransferCount", ctypes.c_ulonglong),
            ]

        class _BasicLimitInformation(ctypes.Structure):
            _fields_ = [
                ("PerProcessUserTimeLimit", ctypes.c_longlong),
                ("PerJobUserTimeLimit", ctypes.c_longlong),
                ("LimitFlags", wintypes.DWORD),
                ("MinimumWorkingSetSize", ctypes.c_size_t),
                ("MaximumWorkingSetSize", ctypes.c_size_t),
                ("ActiveProcessLimit", wintypes.DWORD),
                ("Affinity", ctypes.c_size_t),
                ("PriorityClass", wintypes.DWORD),
                ("SchedulingClass", wintypes.DWORD),
            ]

        class _ExtendedLimitInformation(ctypes.Structure):
            _fields_ = [
                ("BasicLimitInformation", _BasicLimitInformation),
                ("IoInfo", _IoCounters),
                ("ProcessMemoryLimit", ctypes.c_size_t),
                ("JobMemoryLimit", ctypes.c_size_t),
                ("PeakProcessMemoryUsed", ctypes.c_size_t),
                ("PeakJobMemoryUsed", ctypes.c_size_t),
            ]

        kernel32 = ctypes.WinDLL("kernel32", use_last_error=True)
        kernel32.CreateJobObjectW.argtypes = (wintypes.LPVOID, wintypes.LPCWSTR)
        kernel32.CreateJobObjectW.restype = wintypes.HANDLE
        kernel32.SetInformationJobObject.argtypes = (
            wintypes.HANDLE,
            ctypes.c_int,
            wintypes.LPVOID,
            wintypes.DWORD,
        )
        kernel32.SetInformationJobObject.restype = wintypes.BOOL
        kernel32.AssignProcessToJobObject.argtypes = (
            wintypes.HANDLE,
            wintypes.HANDLE,
        )
        kernel32.AssignProcessToJobObject.restype = wintypes.BOOL
        kernel32.CloseHandle.argtypes = (wintypes.HANDLE,)
        kernel32.CloseHandle.restype = wintypes.BOOL

        handle = kernel32.CreateJobObjectW(None, None)
        if not handle:
            raise OSError(ctypes.get_last_error(), "CreateJobObjectW failed")
        self._kernel32 = kernel32
        self._handle = int(handle)
        try:
            information = _ExtendedLimitInformation()
            information.BasicLimitInformation.LimitFlags = 0x00002000
            if not kernel32.SetInformationJobObject(
                handle,
                9,
                ctypes.byref(information),
                ctypes.sizeof(information),
            ):
                raise OSError(ctypes.get_last_error(), "SetInformationJobObject failed")
            process_handle = getattr(process, "_handle", None)
            if process_handle is None or not kernel32.AssignProcessToJobObject(
                handle, wintypes.HANDLE(int(process_handle))
            ):
                raise OSError(
                    ctypes.get_last_error(), "AssignProcessToJobObject failed"
                )
        except BaseException:
            self.close()
            raise

    def close(self) -> None:
        if self._handle is None:
            return
        handle, self._handle = self._handle, None
        self._kernel32.CloseHandle(handle)


def _terminate_parser_process(
    process: multiprocessing.Process,
    job: _WindowsParserJob | None,
    *,
    grace_seconds: float = 2.0,
) -> None:
    if os.name == "nt":
        if job is not None:
            job.close()
        elif process.pid is not None and process.is_alive():
            _terminate_windows_process_tree(process.pid)
    elif process.pid is not None and process.is_alive():
        try:
            os.killpg(process.pid, signal.SIGTERM)
        except (OSError, ProcessLookupError):
            process.terminate()
    process.join(timeout=grace_seconds)
    if process.is_alive():
        if os.name != "nt" and process.pid is not None:
            try:
                os.killpg(process.pid, signal.SIGKILL)
            except (OSError, ProcessLookupError):
                process.kill()
        else:
            process.kill()
        process.join(timeout=grace_seconds)
    if process.is_alive():
        raise ParserProcessError(f"parser process {process.pid} could not be reaped")


def _run_parser_isolated(
    path: Path,
    manifest: SourceManifest,
    docling_path: Path | None,
    *,
    timeout_seconds: float,
    heartbeat_interval_seconds: float,
    result_max_bytes: int,
    temp_dir: Path,
    progress: Callable[[dict[str, Any]], None] | None = None,
    should_stop: Callable[[], bool] | None = None,
    parser: Callable[[Path, SourceManifest, Path | None], _Normalized] | None = None,
) -> _Normalized:
    if timeout_seconds <= 0 or heartbeat_interval_seconds <= 0:
        raise ValueError("parser timeout and heartbeat interval must be positive")
    if heartbeat_interval_seconds >= timeout_seconds:
        raise ValueError("parser heartbeat interval must be less than timeout")
    if isinstance(result_max_bytes, bool) or result_max_bytes <= 0:
        raise ValueError("parser result_max_bytes must be positive")
    parser = _normalize_source if parser is None else parser
    temp_dir.mkdir(parents=True, exist_ok=True)
    descriptor, result_name = tempfile.mkstemp(
        prefix=".parser-result-", suffix=".json", dir=temp_dir
    )
    os.close(descriptor)
    result_path = Path(result_name)
    result_path.unlink(missing_ok=True)
    context = multiprocessing.get_context("spawn")
    start_event = context.Event()
    parent_liveness_reader, parent_liveness_writer = context.Pipe(duplex=False)
    process = context.Process(
        target=_parser_process_entry,
        args=(
            start_event,
            str(result_path),
            str(path),
            manifest.to_dict(),
            str(docling_path) if docling_path is not None else None,
            result_max_bytes,
            parent_liveness_reader,
            parser,
        ),
        name="source-catalog-document-parser",
        daemon=False,
    )
    job: _WindowsParserJob | None = None
    ownership_mode = "process_group"
    started_at = time.monotonic()
    try:
        if os.name == "nt":
            # Spawn children re-read the environment at interpreter start;
            # force UTF-8 so a failing child's stderr (which may carry
            # non-ASCII paths) cannot break UTF-8-decoding collectors
            # (pytest capture).
            os.environ.setdefault("PYTHONUTF8", "1")
        process.start()
        parent_liveness_reader.close()
        if os.name == "nt":
            try:
                job = _WindowsParserJob(process)
                ownership_mode = "windows_job"
            except OSError:
                # Some test/desktop hosts place this process in a restrictive job.
                # The child-side parent monitor remains the ownership fallback.
                job = None
                ownership_mode = "parent_monitor"
        start_event.set()

        def emit_progress() -> None:
            if progress is not None:
                progress(
                    {
                        "detail": "parser_alive",
                        "parser_pid": process.pid,
                        "parser_elapsed_seconds": round(
                            time.monotonic() - started_at, 3
                        ),
                        "parser_timeout_seconds": timeout_seconds,
                        "parser_ownership": ownership_mode,
                    }
                )

        emit_progress()
        while process.is_alive():
            if should_stop is not None and should_stop():
                _terminate_parser_process(process, job)
                job = None
                raise NormalizationCancelledError(
                    f"parser process {process.pid} cancelled by stop request"
                )
            elapsed = time.monotonic() - started_at
            remaining = timeout_seconds - elapsed
            if remaining <= 0:
                _terminate_parser_process(process, job)
                job = None
                raise NormalizationTimeoutError(
                    f"parser process {process.pid} exceeded {timeout_seconds}s for {path}"
                )
            process.join(timeout=min(heartbeat_interval_seconds, remaining))
            if process.is_alive():
                emit_progress()
        process.join()
        if job is not None:
            job.close()
            job = None
        if not result_path.is_file():
            raise ParserProcessError(
                f"parser process {process.pid} exited {process.exitcode} without a result"
            )
        result_size = result_path.stat().st_size
        if result_size > result_max_bytes + 65_536:
            raise ParserResultTooLargeError(
                f"parser envelope is {result_size} bytes; limit is {result_max_bytes}"
            )
        try:
            envelope = json.loads(result_path.read_text(encoding="utf-8"))
        except (OSError, UnicodeDecodeError, json.JSONDecodeError) as exc:
            raise ParserResultProtocolError(
                f"parser result envelope is unreadable: {type(exc).__name__}: {exc}"
            ) from exc
        if not isinstance(envelope, dict) or envelope.get("schema_version") != "1.0":
            raise ParserResultProtocolError("parser result envelope schema is invalid")
        if envelope.get("ok") is not True:
            error_type = str(envelope.get("error_type") or "ParserProcessError")
            message = str(envelope.get("error") or "parser child failed")
            if error_type == "ParserResultTooLargeError":
                raise ParserResultTooLargeError(message)
            if error_type in _UNSUPPORTED_CHILD_ERROR_TYPES or (
                error_type == "ValueError"
                and message in _UNSUPPORTED_VALUE_ERROR_MESSAGES
            ):
                raise UnsupportedDocumentError(f"{error_type}: {message}")
            raise ParserProcessError(f"{error_type}: {message}")
        return _normalized_from_payload(
            envelope.get("result"), expected_source_id=manifest.source_id
        )
    finally:
        start_event.set()
        parent_liveness_writer.close()
        parent_liveness_reader.close()
        if process.pid is not None and process.is_alive():
            _terminate_parser_process(process, job)
            job = None
        if job is not None:
            job.close()
        result_path.unlink(missing_ok=True)


def _nfc_lf(text: str) -> str:
    return unicodedata.normalize("NFC", text.replace("\r\n", "\n").replace("\r", "\n"))


def compute_text_fingerprint(text: str) -> str | None:
    """SHA-256 of normalized text, or None for empty/whitespace-only text.

    Normalization is NFC followed by collapsing every run of whitespace to a
    single space and stripping the ends. Two documents whose extracted text is
    identical after this normalization share a fingerprint regardless of
    byte-level differences (re-encoding, watermarking, line-ending differences),
    which is the basis for semantic (non-exact) duplicate detection. Empty or
    whitespace-only text yields None so unreadable or scanned files are excluded
    from semantic grouping.
    """
    if not isinstance(text, str):
        raise TypeError("text must be a string")
    collapsed = " ".join(unicodedata.normalize("NFC", text).split())
    if not collapsed:
        return None
    return hashlib.sha256(collapsed.encode("utf-8")).hexdigest()


def _decode(data: bytes) -> tuple[str, bool]:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return _nfc_lf(data.decode(encoding)), encoding not in {
                "utf-8-sig",
                "utf-8",
            }
        except UnicodeDecodeError:
            continue
    return _nfc_lf(data.decode("latin-1", errors="replace")), True


def _paragraphs(text: str) -> list[str]:
    values = [item.strip() for item in _SENTENCE_BREAK_RE.split(_nfc_lf(text))]
    return [item for item in values if item]


def _parser_result(
    *,
    source_id: str,
    raw_text: str,
    parser_name: str,
    parser_version: str,
    page_number: int | None = None,
    paragraph_index: int | None = None,
    table_index: int | None = None,
    structured_value: Any = None,
    flags: Iterable[QualityFlag | str] = (),
) -> ParserResult:
    normalized_flags = tuple(flags)
    return ParserResult(
        source_id=source_id,
        coordinates=EvidenceCoordinates(
            page_number=page_number,
            paragraph_index=paragraph_index,
            table_index=table_index,
        ),
        raw_text=_nfc_lf(raw_text),
        structured_value=structured_value,
        parser_name=parser_name,
        parser_version=parser_version,
        parse_status=ParseStatus.PARTIAL if normalized_flags else ParseStatus.PARSED,
        quality_flags=normalized_flags,
    )


def _text_markdown(
    path: Path, source_id: str, *, parser_name: str, parser_version: str
) -> _Normalized:
    text, repaired = _decode(path.read_bytes())
    values = _paragraphs(text)
    if not values:
        raise ValueError("empty text output")
    flags = (QualityFlag.ENCODING_REPAIRED,) if repaired else ()
    results = tuple(
        _parser_result(
            source_id=source_id,
            raw_text=value,
            paragraph_index=index,
            parser_name=parser_name,
            parser_version=parser_version,
            flags=flags,
        )
        for index, value in enumerate(values)
    )
    parts = []
    for result in results:
        parts.extend(
            (
                f"<!-- locator: {result.coordinates.locator()} -->",
                result.raw_text or "",
                "",
            )
        )
    return _Normalized(
        body="\n".join(parts).rstrip() + "\n",
        parser_results=results,
        parser_name=parser_name,
        parser_version=parser_version,
        status="partial" if flags else "completed",
        quality_flags=tuple(item.value for item in flags),
    )


def _bbox(value: Any) -> tuple[float, float, float, float] | None:
    try:
        coordinates = tuple(float(item) for item in value)
    except (TypeError, ValueError):
        return None
    if len(coordinates) != 4 or not all(math.isfinite(item) for item in coordinates):
        return None
    x0, y0, x1, y1 = coordinates
    if x1 <= x0 or y1 <= y0:
        return None
    return x0, y0, x1, y1


def _bbox_intersects(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> bool:
    return (
        left[0] < right[2]
        and left[2] > right[0]
        and left[1] < right[3]
        and left[3] > right[1]
    )


def _pymupdf_page_snapshots(document: Iterable[Any]) -> tuple[dict[str, Any], ...]:
    """Extract strict page snapshots without inventing missing table coordinates."""

    snapshots: list[dict[str, Any]] = []
    for page_number, page in enumerate(document, start=1):
        layout_ambiguous = False
        tables: list[dict[str, Any]] = []
        table_bboxes: list[tuple[float, float, float, float]] = []
        find_tables = getattr(page, "find_tables", None)
        if not callable(find_tables):
            layout_ambiguous = True
        else:
            try:
                found_tables = tuple(getattr(find_tables(), "tables", ()))
            except Exception:
                found_tables = ()
                layout_ambiguous = True
            for table in found_tables:
                try:
                    table_bbox = _bbox(getattr(table, "bbox", None))
                    data = tuple(
                        tuple(
                            _nfc_lf(value) if isinstance(value, str) else value
                            for value in row
                        )
                        for row in table.extract()
                    )
                    rows = int(table.row_count)
                    cols = int(table.col_count)
                    markdown_method = getattr(table, "to_markdown", None)
                    if table_bbox is None or not callable(markdown_method):
                        raise ValueError(
                            "table geometry or Markdown API is unavailable"
                        )
                    markdown = _nfc_lf(str(markdown_method())).strip()
                    if not markdown:
                        raise ValueError("table Markdown is empty")
                    tables.append(
                        {
                            "markdown": markdown,
                            "rows": rows,
                            "cols": cols,
                            "data": data,
                        }
                    )
                    table_bboxes.append(table_bbox)
                except Exception:
                    layout_ambiguous = True
        try:
            blocks = page.get_text("blocks", sort=True)
            narrative: list[str] = []
            for block in blocks:
                if not isinstance(block, (list, tuple)) or len(block) < 7:
                    layout_ambiguous = True
                    continue
                if block[6] != 0:
                    continue
                block_bbox = _bbox(block[:4])
                if block_bbox is None:
                    layout_ambiguous = True
                    continue
                if any(
                    _bbox_intersects(block_bbox, table_bbox)
                    for table_bbox in table_bboxes
                ):
                    continue
                text = _nfc_lf(str(block[4])).strip()
                if text:
                    narrative.append(text)
            page_text = "\n\n".join(narrative)
            error = None
        except Exception as exc:
            page_text = ""
            tables = []
            error = f"{type(exc).__name__}: {str(exc)[:500]}"
        snapshots.append(
            {
                "page_number": page_number,
                "text": page_text,
                "tables": tuple(tables),
                "quality_score": 1.0,
                "ocr_used": False,
                "ocr_confidence": None,
                "layout_ambiguous": layout_ambiguous,
                "encoding_repaired": False,
                "error": error,
            }
        )
    return tuple(snapshots)


def _render_page_aware_markdown(result: PageAwarePDFResult) -> str:
    body: list[str] = []
    for page_number in range(1, result.page_count + 1):
        body.extend((f"## Page {page_number}", ""))
        page_results = tuple(
            item
            for item in result.parser_results
            if item.coordinates.page_number == page_number
        )
        for item in page_results:
            body.extend((f"<!-- locator: {item.coordinates.locator()} -->", ""))
            if item.coordinates.paragraph_index is not None:
                body.extend((item.raw_text or "", ""))
            elif item.coordinates.table_index is not None:
                row = item.coordinates.row_index
                column = item.coordinates.column_index
                value = item.raw_text if item.raw_text is not None else "null"
                body.extend((f"- Table cell [{row}, {column}]: {value}", ""))
            elif QualityFlag.EMPTY_OUTPUT in item.quality_flags:
                body.extend(("_No extractable text or table on this page._", ""))
            else:
                body.extend(("_Page extraction failed._", ""))
    return "\n".join(body).rstrip() + "\n"


def _pdf_markdown(path: Path, manifest: SourceManifest) -> _Normalized:
    import fitz

    parser_version = str(getattr(fitz, "VersionBind", "1.0.0"))
    with fitz.open(str(path)) as document:
        if bool(getattr(document, "needs_pass", False)):
            raise ValueError("PDF is password protected")
        pages = _pymupdf_page_snapshots(document)
    result = adapt_pdf_pages(
        manifest=manifest,
        pages=pages,
        parser_version=parser_version,
    )
    flags = tuple(
        dict.fromkeys(
            flag.value if isinstance(flag, QualityFlag) else str(flag)
            for parser_result in result.parser_results
            for flag in parser_result.quality_flags
        )
    )
    status = (
        "partial"
        if any(
            item.parse_status is not ParseStatus.PARSED
            for item in result.parser_results
        )
        else "completed"
    )
    return _Normalized(
        body=_render_page_aware_markdown(result),
        parser_results=result.parser_results,
        parser_name=PAGE_AWARE_PDF_PARSER_NAME,
        parser_version=parser_version,
        status=status,
        quality_flags=flags,
        page_count=result.page_count,
        first_page_text=_first_page_text(result),
    )


def _first_page_text(result: PageAwarePDFResult) -> str | None:
    """ZR-502: the raw text of page 1 (for homepage identity checks)."""
    page_one = tuple(
        item
        for item in result.parser_results
        if item.coordinates.page_number == 1 and item.raw_text
    )
    if not page_one:
        return None
    return "\n".join(str(item.raw_text) for item in page_one)


def _docling_markdown(path: Path, source_id: str) -> _Normalized:
    from docling_core.types.doc.document import DoclingDocument

    document = DoclingDocument.load_from_json(str(path))
    parser_version = str(getattr(document, "version", "1.0.0"))
    markdown = _nfc_lf(document.export_to_markdown()).strip()
    if not markdown:
        raise ValueError("Docling export is empty")
    pages: dict[int, list[str]] = {}
    for item in document.texts:
        text = _nfc_lf(str(getattr(item, "text", ""))).strip()
        provenance = getattr(item, "prov", None) or []
        page_number = int(provenance[0].page_no) if provenance else 0
        if text and page_number > 0:
            pages.setdefault(page_number, []).append(text)
    results: list[ParserResult] = []
    for page_number, values in sorted(pages.items()):
        results.append(
            _parser_result(
                source_id=source_id,
                raw_text="\n\n".join(values),
                page_number=page_number,
                parser_name="dayu_docling",
                parser_version=parser_version,
                structured_value={"source": "dayu_docling", "page_number": page_number},
            )
        )
    table_index_by_page: dict[int, int] = {}
    locator_lines = ["", "## Locator index", ""]
    for result in results:
        locator_lines.append(f"- `{result.coordinates.locator()}`")
    for table in document.tables:
        provenance = getattr(table, "prov", None) or []
        if not provenance:
            continue
        page_number = int(provenance[0].page_no)
        table_index = table_index_by_page.get(page_number, 0)
        table_index_by_page[page_number] = table_index + 1
        table_text = _nfc_lf(table.export_to_markdown(doc=document)).strip()
        if not table_text:
            continue
        result = _parser_result(
            source_id=source_id,
            raw_text=table_text,
            page_number=page_number,
            table_index=table_index,
            parser_name="dayu_docling",
            parser_version=parser_version,
            structured_value={
                "kind": "table",
                "page_number": page_number,
                "table_index": table_index,
            },
        )
        results.append(result)
        locator_lines.append(f"- `{result.coordinates.locator()}`")
    if not results:
        raise ValueError("Docling provenance contains no page-aware output")
    return _Normalized(
        body=markdown + "\n" + "\n".join(locator_lines) + "\n",
        parser_results=tuple(results),
        parser_name="dayu_docling",
        parser_version=parser_version,
        status="completed",
        quality_flags=(),
    )


def _html_text_markdown(
    text: str, source_id: str, *, format_name: str, repaired: bool
) -> _Normalized:
    from bs4 import BeautifulSoup
    from markdownify import markdownify

    soup = BeautifulSoup(text, "lxml")
    for item in soup(["script", "style", "noscript"]):
        item.decompose()
    markdown = _nfc_lf(markdownify(str(soup), heading_style="ATX")).strip()
    flags = (QualityFlag.ENCODING_REPAIRED,) if repaired else ()
    result = _parser_result(
        source_id=source_id,
        raw_text=soup.get_text("\n", strip=True),
        paragraph_index=0,
        parser_name="html_markdownify",
        parser_version="1.0.0",
        structured_value={"format": format_name},
        flags=flags,
    )
    return _Normalized(
        body=f"<!-- locator: {result.coordinates.locator()} -->\n\n{markdown}\n",
        parser_results=(result,),
        parser_name="html_markdownify",
        parser_version="1.0.0",
        status="partial" if flags else "completed",
        quality_flags=tuple(item.value for item in flags),
    )


def _html_markdown(path: Path, source_id: str) -> _Normalized:
    text, repaired = _decode(path.read_bytes())
    return _html_text_markdown(
        text,
        source_id,
        format_name=path.suffix.lower(),
        repaired=repaired,
    )


def _mht_markdown(path: Path, source_id: str) -> _Normalized:
    message = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
    for part in message.walk():
        if part.get_content_type() == "text/html":
            charset = part.get_content_charset() or "utf-8"
            payload = part.get_payload(decode=True) or b""
            return _html_text_markdown(
                _nfc_lf(payload.decode(charset, errors="replace")),
                source_id,
                format_name=".mht",
                repaired=False,
            )
    return _text_markdown(
        path, source_id, parser_name="mht_text", parser_version="1.0.0"
    )


def _docx_markdown(path: Path, source_id: str) -> _Normalized:
    from docx import Document

    document = Document(path)
    body: list[str] = []
    raw: list[str] = []
    for paragraph in document.paragraphs:
        value = _nfc_lf(paragraph.text).strip()
        if not value:
            continue
        style = str(getattr(paragraph.style, "name", ""))
        if style.startswith("Heading"):
            try:
                level = max(1, min(6, int(style.split()[-1])))
            except ValueError:
                level = 2
            body.append("#" * level + " " + value)
        else:
            body.append(value)
        raw.append(value)
    for table in document.tables:
        rows = [
            [_nfc_lf(cell.text).strip() for cell in row.cells] for row in table.rows
        ]
        if rows:
            body.append(_markdown_table(rows))
            raw.append("\n".join(" | ".join(row) for row in rows))
    text = "\n\n".join(raw)
    if not text:
        raise ValueError("DOCX contains no extractable text")
    result = _parser_result(
        source_id=source_id,
        raw_text=text,
        paragraph_index=0,
        parser_name="python_docx",
        parser_version="1.0.0",
    )
    return _Normalized(
        body=f"<!-- locator: {result.coordinates.locator()} -->\n\n"
        + "\n\n".join(body)
        + "\n",
        parser_results=(result,),
        parser_name="python_docx",
        parser_version="1.0.0",
        status="completed",
        quality_flags=(),
    )


def _escape_cell(value: Any) -> str:
    return (
        str(value if value is not None else "")
        .replace("\n", " ")
        .replace("|", "\\|")
        .strip()
    )


def _markdown_table(rows: list[list[Any]]) -> str:
    width = max((len(row) for row in rows), default=0)
    if width == 0:
        return ""
    normalized = [
        [_escape_cell(row[index] if index < len(row) else "") for index in range(width)]
        for row in rows
    ]
    header = normalized[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _xlsx_markdown(path: Path, source_id: str) -> _Normalized:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    raw_parts: list[str] = []
    results: list[ParserResult] = []
    for sheet_index, sheet in enumerate(workbook.worksheets):
        rows = [list(row) for row in sheet.iter_rows(values_only=True)]
        while rows and not any(value not in (None, "") for value in rows[-1]):
            rows.pop()
        if not rows:
            continue
        table = _markdown_table(rows)
        parts.extend((f"## Sheet: {sheet.title}", "", table, ""))
        raw = "\n".join(
            " | ".join(_escape_cell(value) for value in row) for row in rows
        )
        raw_parts.append(raw)
        results.append(
            _parser_result(
                source_id=source_id,
                raw_text=raw,
                table_index=sheet_index,
                parser_name="openpyxl",
                parser_version=str(openpyxl.__version__),
                structured_value={"sheet": sheet.title},
            )
        )
    workbook.close()
    if not results:
        raise ValueError("XLSX contains no non-empty sheets")
    return _Normalized(
        "\n".join(parts),
        tuple(results),
        "openpyxl",
        str(openpyxl.__version__),
        "completed",
        (),
    )


def _xls_markdown(path: Path, source_id: str) -> _Normalized:
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    parts: list[str] = []
    results: list[ParserResult] = []
    for sheet_index, sheet in enumerate(workbook.sheets()):
        rows = [sheet.row_values(index) for index in range(sheet.nrows)]
        if not rows:
            continue
        parts.extend((f"## Sheet: {sheet.name}", "", _markdown_table(rows), ""))
        raw = "\n".join(
            " | ".join(_escape_cell(value) for value in row) for row in rows
        )
        results.append(
            _parser_result(
                source_id=source_id,
                raw_text=raw,
                table_index=sheet_index,
                parser_name="xlrd",
                parser_version=str(xlrd.__version__),
                structured_value={"sheet": sheet.name},
            )
        )
    workbook.release_resources()
    if not results:
        raise ValueError("XLS contains no non-empty sheets")
    return _Normalized(
        "\n".join(parts), tuple(results), "xlrd", str(xlrd.__version__), "completed", ()
    )


def _pptx_markdown(path: Path, source_id: str) -> _Normalized:
    from pptx import Presentation
    import pptx

    deck = Presentation(path)
    parts: list[str] = []
    results: list[ParserResult] = []
    for slide_index, slide in enumerate(deck.slides):
        values: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                value = _nfc_lf(shape.text).strip()
                if value:
                    values.append(value)
        if not values:
            continue
        page_number = slide_index + 1
        raw = "\n\n".join(values)
        parts.extend(
            (
                f"## Slide {page_number}",
                "",
                f"<!-- locator: loc:v1/page:{page_number} -->",
                "",
                raw,
                "",
            )
        )
        results.append(
            _parser_result(
                source_id=source_id,
                raw_text=raw,
                page_number=page_number,
                parser_name="python_pptx",
                parser_version=str(pptx.__version__),
            )
        )
    if not results:
        raise ValueError("PPTX contains no extractable text")
    return _Normalized(
        "\n".join(parts),
        tuple(results),
        "python_pptx",
        str(pptx.__version__),
        "completed",
        (),
    )


def _doc_markdown(path: Path, source_id: str) -> _Normalized:
    completed = subprocess.run(
        ["antiword", str(path)],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=120,
    )
    if completed.returncode != 0:
        raise ValueError(
            "antiword failed: "
            + completed.stderr.decode("utf-8", errors="replace")[:300]
        )
    text, _ = _decode(completed.stdout)
    temporary = _paragraphs(text)
    if not temporary:
        raise ValueError("antiword returned empty text")
    result = _parser_result(
        source_id=source_id,
        raw_text="\n\n".join(temporary),
        paragraph_index=0,
        parser_name="antiword",
        parser_version="1.0.0",
    )
    return _Normalized(
        f"<!-- locator: {result.coordinates.locator()} -->\n\n{result.raw_text}\n",
        (result,),
        "antiword",
        "1.0.0",
        "completed",
        (),
    )


def _json_xml_markdown(path: Path, source_id: str) -> _Normalized:
    text, repaired = _decode(path.read_bytes())
    language = "json" if path.suffix.lower() == ".json" else "xml"
    flag_values = (QualityFlag.ENCODING_REPAIRED,) if repaired else ()
    result = _parser_result(
        source_id=source_id,
        raw_text=text,
        paragraph_index=0,
        parser_name="structured_text",
        parser_version="1.0.0",
        structured_value={"language": language},
        flags=flag_values,
    )
    return _Normalized(
        f"<!-- locator: {result.coordinates.locator()} -->\n\n```{language}\n{text}\n```\n",
        (result,),
        "structured_text",
        "1.0.0",
        "partial" if repaired else "completed",
        tuple(item.value for item in flag_values),
    )


def _unsupported(path: Path, reason: str) -> _Normalized:
    body = (
        "## Extraction status\n\n"
        "This source was cataloged, but no trustworthy text adapter is available.\n\n"
        f"- Format: `{path.suffix.lower() or '[none]'}`\n"
        f"- Quality flag: `unsupported_format`\n"
        f"- Reason: {html.escape(reason)}\n"
    )
    return _Normalized(
        body,
        (),
        "unsupported_format",
        "1.0.0",
        "unsupported",
        ("unsupported_format",),
        reason,
    )


def _failed(path: Path, reason: str) -> _Normalized:
    body = (
        "## Extraction status\n\n"
        "This source was cataloged, but the isolated parser did not complete.\n\n"
        f"- Format: `{path.suffix.lower() or '[none]'}`\n"
        "- Quality flag: `parser_failed`\n"
        f"- Reason: {html.escape(reason)}\n"
    )
    return _Normalized(
        body,
        (),
        "isolated_parser",
        "1.0.0",
        "failed",
        ("parser_failed",),
        reason,
    )


def _normalize_source(
    path: Path, manifest: SourceManifest, docling_path: Path | None
) -> _Normalized:
    source_id = manifest.source_id
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        if docling_path is not None:
            try:
                return _docling_markdown(docling_path, source_id)
            except Exception:
                pass
        return _pdf_markdown(path, manifest)
    if suffix in {".txt", ".md", ".csv"}:
        return _text_markdown(
            path, source_id, parser_name="plain_text", parser_version="1.0.0"
        )
    if suffix in {".html", ".htm"}:
        return _html_markdown(path, source_id)
    if suffix == ".mht":
        return _mht_markdown(path, source_id)
    if suffix == ".docx":
        return _docx_markdown(path, source_id)
    if suffix == ".doc":
        return _doc_markdown(path, source_id)
    if suffix == ".xlsx":
        return _xlsx_markdown(path, source_id)
    if suffix == ".xls":
        return _xls_markdown(path, source_id)
    if suffix == ".pptx":
        return _pptx_markdown(path, source_id)
    if suffix in {".json", ".xml", ".xsd"}:
        return _json_xml_markdown(path, source_id)
    return _unsupported(path, "No audited parser is installed for this format")


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _atomic_write(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(path.name + f".{os.getpid()}.tmp")
    try:
        temporary.write_text(text, encoding="utf-8", newline="\n")
        os.replace(temporary, path)
    finally:
        temporary.unlink(missing_ok=True)


def _frontmatter(document: Any, normalized: _Normalized) -> str:
    # ZR-502: homepage identity verification — first-page text against the
    # declared title/publisher; contradiction flags the artifact for review.
    from .homepage_identity import (
        assess_homepage_identity,
        homepage_identity_quality_flag,
    )
    # ZR-503: multi-entity attribution guard — full text against the
    # declared entity; multi-entity content flags for attribution, never a
    # silent single-entity pass.
    from .entity_detection import (
        detect_entities,
        multi_entity_quality_flag,
    )
    # ZR-506: section / chunk / fact assertion — structural layer over the
    # flat locator stream (ZR-504/505).
    from .section_chunk_fact import (
        chunk_spans,
        content_line_count,
        detect_sections,
        extract_facts,
    )

    # document may be a sqlite3.Row (normalize_catalog) or a plain dict
    # (tests/fixtures): .get only exists on the dict, index access works on
    # both but only when the key is present.
    if isinstance(document, dict):
        metadata = document.get("metadata_json") or {}
    else:
        # sqlite3.Row: metadata_json is a JSON string column.
        metadata = (
            json.loads(document["metadata_json"]) if document["metadata_json"] else {}
        )
    inner = metadata.get("acquisition") or metadata.get("dayu_meta") or {}
    identity = assess_homepage_identity(
        normalized.first_page_text,
        title=document["title"],
        publisher=inner.get("publisher"),
    )
    identity_verdict = str(identity["verdict"])
    identity_flag = homepage_identity_quality_flag(identity_verdict)
    flags = list(normalized.quality_flags)
    if identity_flag and identity_flag not in flags:
        flags.append(identity_flag)
    entity_detection = detect_entities(
        normalized.body,
        declared_entity=inner.get("canonical_entity_id") or inner.get("display_name"),
        declared_security_ids=tuple(inner.get("security_ids") or ()),
    )
    entity_verdict = str(entity_detection["verdict"])
    entity_flag = multi_entity_quality_flag(entity_verdict)
    if entity_flag and entity_flag not in flags:
        flags.append(entity_flag)
    # ZR-506: structural assertions over the rendered body — sections
    # (chapter headings), chunks (section line ranges), facts
    # ("指标名：数字+单位").  Empty results stay honest (never fabricated).
    structure_sections = detect_sections(normalized.body)
    structure_chunks = chunk_spans(
        content_line_count(normalized.body), structure_sections
    )
    structure_facts = extract_facts(normalized.body)
    document_structure = {
        "schema_version": "1.0",
        "sections": structure_sections,
        "chunk_count": len(structure_chunks),
        "facts": structure_facts,
    }
    payload = {
        "schema_version": "1.0.0",
        "artifact_role": "normalized",
        "document_id": document["document_id"],
        "source_id": document["primary_source_id"],
        "source_sha256": document["content_sha256"],
        "title": document["title"],
        "document_kind": document["document_kind"],
        "published_date": document["published_date"],
        "normalization_status": normalized.status,
        "parser_name": normalized.parser_name,
        "parser_version": normalized.parser_version,
        "quality_flags": flags,
        # ZR-501: page count from the page-aware parser (None stays honest).
        "page_count": normalized.page_count,
        # ZR-502: homepage identity verdict (consistent / contradiction /
        # unverifiable) with evidence; never a fabricated pass.
        "homepage_identity": identity,
        # ZR-503: multi-entity detection verdict (single / multi_entity /
        # unverifiable) with extracted company-name phrases; multi_entity
        # flags for attribution (fail-closed, zero hardcoded names).
        "detected_entities": entity_detection,
        # ZR-506: structural assertions — sections, chunk_count, facts.
        "document_structure": document_structure,
    }
    return (
        "---\n"
        + yaml.safe_dump(payload, allow_unicode=True, sort_keys=False).rstrip()
        + "\n---\n\n"
    )


def normalize_catalog(
    config: CatalogConfig,
    store: CatalogStore,
    *,
    limit: int | None = None,
    force: bool = False,
    progress: Callable[..., None] | None = None,
    should_stop: Callable[[], bool] | None = None,
    parser_timeout_seconds: float = 3600,
    parser_heartbeat_interval_seconds: float = 15,
    parser_result_max_bytes: int = 268_435_456,
    retry_limit: int = 3,
    retry_backoff_seconds: int = 900,
) -> ProcessingReport:
    if limit is not None and limit <= 0:
        raise ValueError("limit must be positive")
    if retry_limit < 1:
        raise ValueError("retry_limit must be >= 1")
    if retry_backoff_seconds < 0:
        raise ValueError("retry_backoff_seconds must be >= 0")
    now_epoch = time.time()
    sql = """SELECT d.*,s.content_sha256,s.byte_size,s.mime_type,
        existing.metadata_json AS normalization_metadata_json
        FROM documents d JOIN sources s ON s.source_id=d.primary_source_id
        LEFT JOIN artifacts existing ON existing.document_id=d.document_id
        AND existing.artifact_role='normalized' AND existing.generator_name=?
        AND existing.generator_version=?"""
    params: tuple[Any, ...] = (_NORMALIZER_NAME, NORMALIZER_VERSION)
    # FC-906-c: a document without an active original_primary location cannot
    # be parsed — it used to sit at the queue head forever and fail with no
    # artifact row and no last_failed diagnostic (production: 9506/23521 docs).
    # Exclude them from the queue (in force and non-force runs alike) so real
    # documents behind them are never starved.
    sql += """ WHERE EXISTS (
            SELECT 1 FROM locations lp JOIN roots rp ON rp.root_id=lp.root_id
            WHERE lp.document_id=d.document_id AND lp.role='original_primary'
              AND lp.location_status='active')"""
    if not force:
        sql += """ AND (existing.artifact_id IS NULL OR (
            existing.status='failed'
            AND COALESCE(json_extract(existing.metadata_json,'$.terminal'),0)=0
            AND COALESCE(json_extract(existing.metadata_json,'$.next_retry_epoch'),0)<=?)
        )"""
        params += (now_epoch,)
    sql += f" ORDER BY {processing_priority_sql('d')}, d.document_id"
    if limit is not None:
        sql += " LIMIT ?"
        params += (limit,)
    documents = store.fetchall(sql, params)
    completed = skipped = partial = unsupported = failed = 0
    failure_reasons: dict[str, int] = {}
    last_failure_code: str | None = None
    last_failed_document_id: str | None = None
    last_failed_path: str | None = None
    for document_index, document in enumerate(documents, start=1):
        if should_stop is not None and should_stop():
            partial += max(0, len(documents) - document_index + 1)
            break
        source_id = document["primary_source_id"]
        locations = store.fetchall(
            """SELECT l.*,r.path AS root_path,r.priority FROM locations l JOIN roots r ON r.root_id=l.root_id
            WHERE l.document_id=? AND l.location_status='active' ORDER BY r.priority,l.relative_path""",
            (document["document_id"],),
        )
        primary = next(
            (
                item
                for item in locations
                if item["role"] == "original_primary" and item["source_id"] == source_id
            ),
            None,
        )
        if primary is None:
            # FC-906-c defense in depth: a document slipping past the queue
            # filter must be VISIBLE in the report (reason + id + path), not a
            # silent failure count.
            failed += 1
            last_failure_code = "no_active_primary_location"
            last_failed_document_id = document["document_id"]
            last_failed_path = ""
            failure_reasons[last_failure_code] = (
                failure_reasons.get(last_failure_code, 0) + 1
            )
            continue
        source_path = Path(primary["absolute_path"])
        if progress is not None:
            progress(
                current_path=str(source_path.resolve(strict=False)),
                current=document_index,
                total=len(documents),
                detail="extracting Markdown",
            )
        manifest = SourceManifest.from_dict(json.loads(primary["manifest_json"]))
        docling_path: Path | None = None
        metadata = json.loads(document["metadata_json"])
        dayu_meta = metadata.get("dayu_meta") or {}
        expected_pdf_sha = dayu_meta.get("pdf_sha256")
        if (
            source_path.suffix.lower() == ".pdf"
            and expected_pdf_sha == document["content_sha256"]
        ):
            sidecar = next(
                (item for item in locations if item["role"] == "processed_docling"),
                None,
            )
            if sidecar is not None:
                possible = Path(sidecar["absolute_path"])
                if possible.is_file():
                    docling_path = possible

        def parser_progress(details: dict[str, Any]) -> None:
            if progress is not None:
                progress(
                    current_path=str(source_path.resolve(strict=False)),
                    current=document_index,
                    total=len(documents),
                    **details,
                )

        failure_metadata: dict[str, Any] | None = None
        try:
            normalized = _run_parser_isolated(
                source_path,
                manifest,
                docling_path,
                timeout_seconds=parser_timeout_seconds,
                heartbeat_interval_seconds=parser_heartbeat_interval_seconds,
                result_max_bytes=parser_result_max_bytes,
                temp_dir=config.catalog_dir / "parser_tmp",
                progress=parser_progress,
                should_stop=should_stop,
            )
            bundle = IngestService(root=Path(primary["root_path"])).ingest(
                manifest=manifest,
                parser_results=normalized.parser_results,
            )
        except NormalizationCancelledError:
            partial += max(0, len(documents) - document_index + 1)
            break
        except UnsupportedDocumentError as exc:
            normalized = _unsupported(source_path, str(exc)[:500])
            bundle = IngestService(root=Path(primary["root_path"])).ingest(
                manifest=manifest, parser_results=()
            )
        except Exception as exc:
            existing_metadata = json.loads(
                document["normalization_metadata_json"] or "{}"
            )
            next_attempt = int(existing_metadata.get("attempt_count") or 0) + 1
            terminal = next_attempt >= retry_limit
            error_code = (
                DOCUMENT_PARSE_TIMEOUT_CODE
                if isinstance(exc, NormalizationTimeoutError)
                else type(exc).__name__
            )
            error_text = f"{error_code}: {str(exc)[:500]}"
            normalized = _failed(source_path, error_text)
            bundle = IngestService(root=Path(primary["root_path"])).ingest(
                manifest=manifest, parser_results=()
            )
            failure_metadata = {
                "attempt_count": next_attempt,
                "terminal": terminal,
                "terminal_reason": (
                    f"retry_exhausted:{error_code}" if terminal else None
                ),
                "next_retry_epoch": (
                    None if terminal else now_epoch + retry_backoff_seconds
                ),
                "error_code": error_code,
            }
            failed += 1
            last_failure_code = error_code
            last_failed_document_id = document["document_id"]
            last_failed_path = str(source_path.resolve(strict=False))
            failure_reasons[error_code] = failure_reasons.get(error_code, 0) + 1
        raw_text = "\n\n".join(
            result.raw_text or "" for result in normalized.parser_results
        )
        text_fingerprint = compute_text_fingerprint(raw_text)
        output_path = (
            config.derived_dir
            / document["content_sha256"][:2]
            / document["content_sha256"]
            / "normalized.md"
        )
        content = (
            _frontmatter(document, normalized)
            + f"# {document['title']}\n\n"
            + normalized.body
        )
        _atomic_write(output_path, content)
        artifact_hash = _sha256_file(output_path)
        artifact_id = (
            "urn:company-wiki:artifact:sha256:"
            + hashlib.sha256(
                (
                    document["document_id"] + "\0normalized\0" + NORMALIZER_VERSION
                ).encode("utf-8")
            ).hexdigest()
        )
        with store.transaction() as connection:
            connection.execute(
                "DELETE FROM evidence_spans WHERE document_id=?",
                (document["document_id"],),
            )
            for span in bundle.evidence_spans:
                data = span.to_dict()
                coordinates = data["coordinates"]
                connection.execute(
                    """INSERT INTO evidence_spans(span_id,document_id,source_id,locator,page_number,
                    paragraph_index,table_index,raw_text,span_json,parser_name,parser_version,parse_status)
                    VALUES(?,?,?,?,?,?,?,?,?,?,?,?)""",
                    (
                        span.span_id,
                        document["document_id"],
                        span.source_id,
                        span.locator,
                        coordinates["page_number"],
                        coordinates["paragraph_index"],
                        coordinates["table_index"],
                        span.raw_text,
                        span.canonical_json(),
                        span.parser_name,
                        span.parser_version,
                        span.parse_status.value,
                    ),
                )
            connection.execute(
                """INSERT INTO artifacts(artifact_id,document_id,source_id,artifact_role,path,content_sha256,
                byte_size,mime_type,generator_name,generator_version,status,error,
                schema_version,source_sha256,metadata_json,created_at)
                VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,strftime('%Y-%m-%dT%H:%M:%SZ','now'))
                ON CONFLICT(document_id,artifact_role,generator_name,generator_version) DO UPDATE SET
                path=excluded.path,content_sha256=excluded.content_sha256,byte_size=excluded.byte_size,
                status=excluded.status,error=excluded.error,
                schema_version=excluded.schema_version,source_sha256=excluded.source_sha256,
                metadata_json=excluded.metadata_json,created_at=excluded.created_at""",
                (
                    artifact_id,
                    document["document_id"],
                    source_id,
                    "normalized",
                    str(output_path.resolve()),
                    artifact_hash,
                    output_path.stat().st_size,
                    "text/markdown",
                    _NORMALIZER_NAME,
                    NORMALIZER_VERSION,
                    normalized.status,
                    normalized.error,
                    ARTIFACT_HANDLE_SCHEMA_VERSION,
                    str(document["content_sha256"] or ""),
                    canonical_json(
                        {
                            "schema_version": ARTIFACT_HANDLE_SCHEMA_VERSION,
                            "parser_name": normalized.parser_name,
                            "parser_version": normalized.parser_version,
                            "quality_flags": list(normalized.quality_flags),
                            "span_count": len(bundle.evidence_spans),
                            **(failure_metadata or {}),
                        }
                    ),
                ),
            )
            connection.execute(
                "UPDATE documents SET text_fingerprint=? WHERE document_id=?",
                (text_fingerprint, document["document_id"]),
            )
        if normalized.status == "completed":
            completed += 1
        elif normalized.status == "partial":
            partial += 1
        elif normalized.status == "unsupported":
            unsupported += 1
    return ProcessingReport(
        "normalize",
        completed,
        skipped,
        partial,
        unsupported,
        failed,
        terminal_reasons=failure_reasons or None,
        last_failure_code=last_failure_code,
        last_failed_document_id=last_failed_document_id,
        last_failed_path=last_failed_path,
    )


def backfill_text_fingerprints(
    config: CatalogConfig,
    store: CatalogStore,
    *,
    limit: int | None = None,
    progress: Callable[..., None] | None = None,
    should_stop: Callable[[], bool] | None = None,
    retry_limit: int = 3,
    retry_backoff_seconds: int = 900,
    now_epoch: float | None = None,
    parser_timeout_seconds: float = 3600,
    parser_heartbeat_interval_seconds: float = 15,
    parser_result_max_bytes: int = 268_435_456,
) -> ProcessingReport:
    """Compute and persist ``text_fingerprint`` via the persistent state machine.

    CW-2.28 §12.3 / §12.4.3.3. Dispatch reads from ``document_fingerprint_state``
    (pending + due retryable_failed) instead of re-selecting every NULL row, so
    terminal documents are never re-attempted and retryable failures respect a
    bounded retry/backoff. Each document's outcome is written atomically
    (``documents.text_fingerprint`` + state row in one transaction), so a crash
    cannot split the two writes. No normalized/summary artifacts are written.

    Outcome classification:
      * parseable non-empty text  → ``completed`` (+ fingerprint);
      * empty/whitespace text     → ``unsupported_terminal`` (reason ``empty_text``);
      * no original_primary location → ``unsupported_terminal`` (``no_original_location``);
      * parser/I-O exception      → ``retryable_failed`` with backoff, or
                                    ``failed_terminal`` (``retry_exhausted:<code>``)
                                    once ``attempt_count`` reaches ``retry_limit``.

    ``should_stop`` is checked before each document; when it returns True the
    current file completes cleanly and the batch stops (partial, not failed).
    """
    if limit is not None and limit <= 0:
        raise ValueError("limit must be positive")
    if retry_limit < 1:
        raise ValueError("retry_limit must be >= 1")
    if retry_backoff_seconds < 0:
        raise ValueError("retry_backoff_seconds must be >= 0")

    import time as _time

    now = now_epoch if now_epoch is not None else _time.time()
    now_iso = _utc_iso(now)

    # Global backlog before the batch (pending + due retryable_failed).
    status_before = store.fingerprint_status(now_iso=now_iso)
    eligible_total = status_before["eligible"]

    batch = store.select_fingerprint_batch(limit=limit, now_iso=now_iso)
    completed = skipped = partial = unsupported = failed = 0
    terminal_reasons: dict[str, int] = {}
    for document_index, document in enumerate(batch, start=1):
        if should_stop is not None and should_stop():
            partial += max(0, len(batch) - document_index + 1)
            break
        document_id = document["document_id"]
        source_id = document["source_id"]
        source_sha256 = document["source_sha256"]
        attempt_count = int(document["attempt_count"])
        locations = store.fetchall(
            """SELECT l.*,r.path AS root_path FROM locations l JOIN roots r ON r.root_id=l.root_id
            WHERE l.document_id=? AND l.location_status='active' AND l.role='original_primary'
            AND l.source_id=? ORDER BY r.priority,l.relative_path""",
            (document_id, source_id),
        )
        if not locations:
            store.record_fingerprint_outcome(
                document_id=document_id,
                source_id=source_id,
                source_sha256=source_sha256,
                fingerprint=None,
                status="unsupported_terminal",
                attempt_count=attempt_count + 1,
                terminal_reason="no_original_location",
                updated_at=now_iso,
            )
            unsupported += 1
            terminal_reasons["no_original_location"] = (
                terminal_reasons.get("no_original_location", 0) + 1
            )
            continue
        source_path = Path(locations[0]["absolute_path"])
        if progress is not None:
            progress(
                current_path=str(source_path.resolve(strict=False)),
                current=document_index,
                total=len(batch),
                detail="backfilling text fingerprint",
            )
        manifest = SourceManifest.from_dict(json.loads(locations[0]["manifest_json"]))

        def parser_progress(details: dict[str, Any]) -> None:
            if progress is not None:
                progress(
                    current_path=str(source_path.resolve(strict=False)),
                    current=document_index,
                    total=len(batch),
                    **details,
                )

        try:
            normalized = _run_parser_isolated(
                source_path,
                manifest,
                None,
                timeout_seconds=parser_timeout_seconds,
                heartbeat_interval_seconds=parser_heartbeat_interval_seconds,
                result_max_bytes=parser_result_max_bytes,
                temp_dir=config.catalog_dir / "parser_tmp",
                progress=parser_progress,
                should_stop=should_stop,
            )
            raw_text = "\n\n".join(
                result.raw_text or "" for result in normalized.parser_results
            )
            fingerprint = compute_text_fingerprint(raw_text)
        except NormalizationCancelledError:
            partial += max(0, len(batch) - document_index + 1)
            break
        except UnsupportedDocumentError as exc:
            store.record_fingerprint_outcome(
                document_id=document_id,
                source_id=source_id,
                source_sha256=source_sha256,
                fingerprint=None,
                status="unsupported_terminal",
                attempt_count=attempt_count + 1,
                terminal_reason="unsupported_document",
                error_code="unsupported_document",
                error_message=str(exc),
                updated_at=now_iso,
            )
            unsupported += 1
            terminal_reasons["unsupported_document"] = (
                terminal_reasons.get("unsupported_document", 0) + 1
            )
            continue
        except Exception as exc:
            next_attempt = attempt_count + 1
            error_code = (
                DOCUMENT_PARSE_TIMEOUT_CODE
                if isinstance(exc, NormalizationTimeoutError)
                else type(exc).__name__
            )
            if next_attempt >= retry_limit:
                store.record_fingerprint_outcome(
                    document_id=document_id,
                    source_id=source_id,
                    source_sha256=source_sha256,
                    fingerprint=None,
                    status="failed_terminal",
                    attempt_count=next_attempt,
                    terminal_reason=f"retry_exhausted:{error_code}",
                    error_code=error_code,
                    error_message=str(exc),
                    updated_at=now_iso,
                )
                terminal_reasons[f"retry_exhausted:{error_code}"] = (
                    terminal_reasons.get(f"retry_exhausted:{error_code}", 0) + 1
                )
            else:
                next_retry_iso = _utc_iso(now + retry_backoff_seconds)
                store.record_fingerprint_outcome(
                    document_id=document_id,
                    source_id=source_id,
                    source_sha256=source_sha256,
                    fingerprint=None,
                    status="retryable_failed",
                    attempt_count=next_attempt,
                    error_code=error_code,
                    error_message=str(exc),
                    next_retry_at=next_retry_iso,
                    updated_at=now_iso,
                )
            failed += 1
            continue
        if fingerprint is None:
            store.record_fingerprint_outcome(
                document_id=document_id,
                source_id=source_id,
                source_sha256=source_sha256,
                fingerprint=None,
                status="unsupported_terminal",
                attempt_count=attempt_count + 1,
                terminal_reason="empty_text",
                updated_at=now_iso,
            )
            unsupported += 1
            terminal_reasons["empty_text"] = terminal_reasons.get("empty_text", 0) + 1
            continue
        store.record_fingerprint_outcome(
            document_id=document_id,
            source_id=source_id,
            source_sha256=source_sha256,
            fingerprint=fingerprint,
            status="completed",
            attempt_count=attempt_count + 1,
            updated_at=now_iso,
        )
        completed += 1

    status_after = store.fingerprint_status(now_iso=now_iso)
    report = ProcessingReport(
        "backfill_text_fingerprints",
        completed=completed,
        skipped=skipped,
        partial=partial,
        unsupported=unsupported,
        failed=failed,
        eligible=eligible_total,
        terminal_reasons=terminal_reasons if terminal_reasons else None,
        due_retry=status_after["due_retry"],
        terminal=status_after["terminal"],
    )
    return report


__all__ = [
    "backfill_text_fingerprints",
    "compute_text_fingerprint",
    "normalize_catalog",
]
